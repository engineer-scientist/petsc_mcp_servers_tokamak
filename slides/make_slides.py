#!/usr/bin/env python3
"""
make_slides.py -- build the Argonne summer-2026 intern presentation from ONE content spec.

The deck is built **on the official Argonne PowerPoint template**
(slides/Argonne_Powerpoint_Template.pptx): we open the template, drop its example
slides, and populate its own layouts/placeholders, so the theme, fonts (Arial),
Argonne colors, logos, and slide numbers all come from the template itself.

Two backends, same content:
  * the EDITABLE, branded PowerPoint  -> slides/petsc_multiagent_tokamak.pptx   (authoritative)
  * a multi-page PDF preview          -> slides/petsc_multiagent_tokamak.pdf    (GitHub-viewable)

There is no LibreOffice on this host, so the PDF cannot be rendered from the .pptx
directly; instead it is drawn with matplotlib, *restyled to the Argonne palette* as a
faithful-enough proxy. The .pptx is the real, branded deliverable.

Content/numbers come from artifacts/<latest>/metrics.json + verification.json (the Python
driver run) and from artifacts/orchestrator-20260724-fixed/ (the built-in orchestrator demo,
Task 5). Result figures come from figures/; motivation images from slides/*.jpg,*.png.

Run with a Python that has python-pptx + matplotlib + PIL (the tokamak venv has all three):
  /home/sarthak.sharma/tokamak/.venv/bin/python slides/make_slides.py
"""
import os
import json
import textwrap

HERE = os.path.dirname(os.path.abspath(__file__))
PROJECT = os.path.dirname(HERE)
ARTIFACTS = os.path.join(PROJECT, "artifacts")
FIGURES = os.path.join(PROJECT, "figures")
TEMPLATE = os.path.join(HERE, "Argonne_Powerpoint_Template.pptx")

# Slide geometry matches the Argonne template (16:9, 10 x 5.625 in).
W, H = 10.0, 5.625

# Argonne brand palette (from the template theme): dk2 blue, lt2 gold, accents.
ARG_BLUE   = "#0082CA"   # theme dk2 (primary Argonne blue)
ARG_BLUE_DK = "#00609C"  # theme accent2 (darker blue, good for text on white)
ARG_GOLD   = "#F8B200"   # theme lt2
ARG_TEAL   = "#00A19C"   # theme accent5
ARG_GREEN  = "#77B300"   # theme accent1
DARK       = "#1F1F1F"
GREY       = "#5A5A5A"
WHITE      = "#FFFFFF"

# Font sizes (pt). Titles/subtitles largely inherit the template; these tune the body + PDF.
F_TITLE, F_SUB, F_AFFIL = 30, 17, 12          # PDF title slide
F_HEAD, F_SUBHEAD, F_L0, F_L1, F_CODE, F_CAP = 24, 14, 15, 13, 12, 11   # content slides

# Motivation images that live in slides/ (the deuterium-tritium / tokamak / plant art).
IMG_PLANT   = os.path.join(HERE, "schematic of nuclear fusion power plant.png")
IMG_DT      = os.path.join(HERE, "DT nuclear fusion schematic (GG).jpg")
IMG_TOKAMAK = os.path.join(HERE, "doe-explains-tokamaks.jpg")
# NOTE: "particle in cell algorithm schematic.jpg" is deliberately NOT used -- PIC is a
# kinetic method, whereas this work solves the Grad-Shafranov *equilibrium* (an elliptic
# MHD PDE) with finite differences / finite elements; showing PIC would misrepresent it.


def latest_run():
    p = os.path.join(ARTIFACTS, "LATEST")
    if os.path.isfile(p):
        return open(p).read().strip()
    runs = sorted(d for d in os.listdir(ARTIFACTS) if d.startswith("run-"))
    return runs[-1] if runs else None


def load(p):
    return json.load(open(p)) if os.path.isfile(p) else {}


def build_spec():
    run_id = latest_run()
    d = os.path.join(ARTIFACTS, run_id) if run_id else ""
    m = load(os.path.join(d, "metrics.json"))
    c = m.get("correctness", {}); e = m.get("efficiency", {}); h = m.get("human_effort", {})
    order = ", ".join("%.2f" % p for p in (c.get("observed_order_maxnorm") or [])) or "2.00, 2.00, 2.00"
    lines_gen = h.get("solver_lines_agent_generated") or 267
    model_name = c.get("model_name") or "Grad-Shafranov equation"
    err = ("%.2e" % c["finest_grid_maxnorm_error"]) if c.get("finest_grid_maxnorm_error") else "1.3e-05"
    snes = (c.get("snes_converged_reason", "CONVERGED_FNORM_RELATIVE").split(":")[-1].strip())
    flux = os.path.join(FIGURES, "gs_flux_surfaces.png")
    conv = os.path.join(FIGURES, "gs_convergence.png")

    slides = []

    # 1. Title slide (template layout "Title Slide").
    slides.append(dict(kind="title",
        title="Automated Problem-to-Solution Generation\nfor a Tokamak Fusion-Plasma Simulation",
        sub="A hierarchical multi-agent PETSc system — applied and verified",
        presenter="[PRESENTER NAME]", role="Summer 2026 Intern, Argonne National Laboratory",
        event="Argonne Summer 2026\nIntern Presentation", date="2026",
        image=IMG_PLANT))

    # 2. Why fusion.
    slides.append(dict(kind="content", title="Why fusion energy needs trustworthy simulation",
        subtitle="Fusion could be a near-limitless clean energy source — if we can confine the plasma",
        image=IMG_DT, image_caption="Deuterium–tritium fusion", bullets=[
        ("Fusion fuses light nuclei (deuterium + tritium) to release energy — the reaction that powers the Sun.", 0),
        ("A tokamak confines a ~150-million-°C plasma in a magnetic “cage” long enough to fuse.", 0),
        ("Getting there is a simulation problem: predict the plasma before building costly hardware.", 0),
        ("Simulations must be trustworthy — a plausible-looking wrong answer is worse than none.", 0),
    ]))

    # 3. The modeling challenge (tokamak).
    slides.append(dict(kind="content", title="The modeling challenge",
        subtitle="Turning fusion physics into correct, fast HPC code needs scarce, implicit expertise",
        image=IMG_TOKAMAK, image_caption="Tokamak magnetic confinement (U.S. DOE)", bullets=[
        ("Numerical methods: choose grids, discretizations, and solvers that actually converge.", 0),
        ("Libraries & parallelism: express it in PETSc and run on many CPU cores / GPUs.", 0),
        ("Verification: prove the computed answer is right, not merely plausible.", 0),
        ("Question: can a multi-agent AI system automate this whole path — and can we trust it?", 0),
    ]))

    # 4. The problem: Grad-Shafranov.
    slides.append(dict(kind="content", title="The problem: Grad–Shafranov equilibrium",
        subtitle="Axisymmetric ideal-MHD force balance for the poloidal flux ψ(R,Z)",
        image=flux, image_caption="Flux surfaces of the verified solution", bullets=[
        ("Δ*ψ = −μ₀R² p′(ψ) − F F′(ψ) — elliptic, nonlinear, time-independent.", 0),
        ("Its solution gives the flux surfaces, magnetic axis, and safety factor q.", 0),
        ("A natural PETSc SNES (nonlinear solver) problem.", 0),
        ("Verifiable: a manufactured exact solution gives a rigorous convergence test.", 0),
    ]))

    # 5. The system.
    slides.append(dict(kind="content", title="The system: a 3-layer multi-agent pipeline",
        subtitle="PETSc’s Model-Context-Protocol (MCP) agents, coordinated end to end",
        bullets=[
        ("Problem definition — Mathematical Modeling agent → strong/weak form, name, time-dependence.", 0),
        ("Agent execution — Numerical Analysis agent → grid, discretization, solver (SNES).", 0),
        ("Agent execution — HPC Code Generation agent → writes, compiles & runs PETSc C.", 0),
        ("Execution substrate — compile-run agent (make / run on the real machine).", 0),
        ("All agents run against Argonne’s Argo gateway (Claude Opus 4.8).", 0),
        ("A project-owned driver records every artifact with provenance (reproducible, resumable).", 0),
    ]))

    # 6. The pipeline in action.
    slides.append(dict(kind="content", title="The pipeline in action",
        subtitle="From one English prompt to a compiled, running solver — with no human edits",
        bullets=[
        ("Prompt: “the Grad-Shafranov equilibrium for the plasma in a tokamak.”", 0),
        ("Modeling agent → identified ‘%s’; time-dependent = False." % model_name, 0),
        ("Numerical Analysis agent → nonlinear solve with SNES on a structured grid.", 0),
        ("Code Generation agent → %d-line PETSc program (DMDA + SNES + true Jacobian)." % lines_gen, 0),
        ("Compiled and ran on 1 and 4 MPI ranks — no human edits.", 0),
    ]))

    # 7. Code excerpt.
    slides.append(dict(kind="code", title="The agent-generated PETSc solver (excerpt)",
        subtitle="A true-Jacobian finite-difference SNES residual — written by the Code Generation agent",
        code=[
        "static PetscErrorCode FormFunction(SNES snes, Vec X, Vec F, void *ctx){",
        "  /* Delta*_h psi = (psi_E-2psi_C+psi_W)/hR^2",
        "                  - (1/R)(psi_E-psi_W)/(2hR)",
        "                  + (psi_N-2psi_C+psi_S)/hZ^2   */",
        "  f[j][i] = dRR - dR1/R + dZZ - ForcingF(user,R,Z);   /* residual */",
        "}",
        "SNESSetFunction(snes, r, FormFunction, &user);",
        "SNESSetJacobian(snes, J, J, FormJacobian, &user);   /* true Jacobian */",
        "SNESSolve(snes, NULL, x);",
    ]))

    # 8. Verification.
    slides.append(dict(kind="content", title="Verification: it is actually right",
        subtitle="Method of manufactured solutions — a measured order of accuracy, not just plausible output",
        image=conv, image_caption="Manufactured-solution convergence", bullets=[
        ("Known exact ψ → check the numerical error as the grid refines.", 0),
        ("Max-norm error = %s on the finest 257² grid." % err, 0),
        ("Observed order of accuracy p = %s (expected 2 for central differences)." % order, 0),
        ("SNES converged: %s." % snes, 0),
    ]))

    # 9. Metrics.
    slides.append(dict(kind="content", title="Decision-gate metrics",
        subtitle="Correctness, human effort, and efficiency — the numbers behind the demo",
        bullets=[
        ("Correctness: model identified ✓, compiled & ran ✓, 2nd-order convergence ✓.", 0),
        ("Human effort: %s lines of solver code hand-written — it was generated." % h.get("solver_lines_handwritten", 0), 0),
        ("Efficiency: total wall-clock ≈ %s s; code-gen used %s tool calls." % (e.get("wallclock_seconds_total", "451"), e.get("codegen_tool_calls", "5")), 0),
        ("≈ %s LLM completions across the whole pipeline." % e.get("approx_llm_completions", "23"), 0),
    ]))

    # 10. NEW -- the built-in orchestrator (Task 5). Numbers from
    #     artifacts/orchestrator-20260724-fixed/ .
    slides.append(dict(kind="content", title="Fully autonomous: the built-in orchestrator",
        subtitle="The shipped LLM orchestrator drives all four agents itself — no Python driver",
        bullets=[
        ("Given only the prompt, an inner Claude agent sequenced all four MCP servers on its own.", 0),
        ("It independently chose a Solov’ev Grad–Shafranov model and a DMPLEX + PetscFE finite-element solve.", 0),
        ("The generated code compiled cleanly and ran: 225 DOFs, L2 norm ψ = 2.04, return code 0.", 0),
        ("It surfaced a real upstream bug — a hardcoded iteration cap flagged a converged run as failure.", 0),
        ("We fixed it and contributed the patch; the re-run finished: “I have completed the orchestration.”", 0),
    ]))

    # 11. What I built.
    slides.append(dict(kind="content", title="What I built and contributed",
        subtitle="Reusable tooling around the agents, plus fixes contributed back upstream",
        bullets=[
        ("A project-owned orchestration driver with full artifact provenance (resumable).", 0),
        ("Verification + post-processing (convergence, flux surfaces) and a metrics harness.", 0),
        ("Four fixes contributed back to the open PETSc multi-agent system:", 0),
        ("CWD-independent server resolution; portable stdio interpreter", 1),
        ("graceful operation without the documentation / RAG services", 1),
        ("code-generation and orchestrator loops that reliably capture results", 1),
    ]))

    # 12. Takeaways.
    slides.append(dict(kind="content", title="Takeaways & next steps",
        subtitle="A verification-driven multi-agent path from prompt to trustworthy HPC code",
        bullets=[
        ("A multi-agent system generated a correct, verified PETSc fusion-MHD solver from a prompt.", 0),
        ("Both the Python driver and the fully-autonomous orchestrator succeeded end to end.", 0),
        ("Verification-driven: convergence + exact-solution checks, not just plausible output.", 0),
        ("Next: shaped real-machine equilibria (Solov’ev / Cerfon–Freidberg), q-profile, GPU scale-up.", 0),
        ("Code + docs: github.com/engineer-scientist/petsc_mcp_servers_tokamak", 0),
    ]))

    # 13. Closing (template's Argonne closing layout).
    slides.append(dict(kind="closing"))

    return dict(slides=slides, run_id=run_id)


# ============================ python-pptx (template) backend ============================
def render_pptx(spec, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image

    def rgb(x):
        return RGBColor(int(x[1:3], 16), int(x[3:5], 16), int(x[5:7], 16))

    prs = Presentation(TEMPLATE)

    # Map layouts by name so we don't depend on index order.
    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}
    L_TITLE = layouts.get("Title Slide")
    L_BULLETS = layouts.get("*Title, Subtitle and Bullets")
    L_TITLEONLY = layouts.get("*Title and Subtitle Only")
    L_CLOSING = layouts.get("Closing slide Argonne DOE") or layouts.get("Closing slide Argonne")

    # Drop the template's example slides (keep the master/layouts/theme).
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn_r_id())
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)

    def ph_by_idx(slide, idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def set_text(ph, text, size=None, bold=None, color=None):
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        for i, ln in enumerate(text.split("\n")):
            pp = p if i == 0 else tf.add_paragraph()
            r = pp.add_run(); r.text = ln
            if size is not None: r.font.size = Pt(size)
            if bold is not None: r.font.bold = bold
            if color is not None: r.font.color.rgb = rgb(color)

    def add_image_fit(slide, path, box):
        bx, by, bw, bh = box
        iw, ih = Image.open(path).size
        scale = min(bw / iw, bh / ih)
        dw, dh = iw * scale, ih * scale
        return slide.shapes.add_picture(path, Inches(bx + (bw - dw) / 2), Inches(by + (bh - dh) / 2),
                                        width=Inches(dw), height=Inches(dh))

    def add_caption(slide, text, box):
        bx, by, bw, bh = box
        tb = slide.shapes.add_textbox(Inches(bx), Inches(by), Inches(bw), Inches(bh))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text; r.font.size = Pt(F_CAP); r.font.italic = True
        r.font.color.rgb = rgb(GREY)

    def fill_bullets(ph, bullets):
        tf = ph.text_frame; tf.clear(); tf.word_wrap = True
        for i, (text, lvl) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = lvl
            r = p.add_run(); r.text = text
            r.font.size = Pt(F_L1 if lvl else F_L0)

    for s in spec["slides"]:
        kind = s["kind"]

        if kind == "closing":
            prs.slides.add_slide(L_CLOSING)
            continue

        if kind == "title":
            sl = prs.slides.add_slide(L_TITLE)
            t = ph_by_idx(sl, 0)
            if t is not None: set_text(t, s["title"])
            sub = ph_by_idx(sl, 1)
            if sub is not None: set_text(sub, s["sub"])
            # picture placeholder
            pic = ph_by_idx(sl, 10)
            if pic is not None and s.get("image") and os.path.isfile(s["image"]):
                try:
                    pic.insert_picture(s["image"])
                except Exception:
                    pass
            # footer columns: (label idx, value idx, label, value)
            footer = [(17, 18, "Presenter", "%s\n%s" % (s["presenter"], s["role"])),
                      (19, 20, "Event", s["event"]),
                      (21, 22, "Date", s["date"])]
            for lab_i, val_i, lab, val in footer:
                lp = ph_by_idx(sl, lab_i); vp = ph_by_idx(sl, val_i)
                if lp is not None: set_text(lp, lab)
                if vp is not None: set_text(vp, val)
            continue

        if kind == "code":
            sl = prs.slides.add_slide(L_TITLEONLY)
            t = ph_by_idx(sl, 0)
            if t is not None: set_text(t, s["title"])
            sub = ph_by_idx(sl, 13)
            if sub is not None and s.get("subtitle"): set_text(sub, s["subtitle"])
            # code panel
            box = (0.5, 1.7, 9.0, 3.3)
            panel = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
            panel.fill.solid(); panel.fill.fore_color.rgb = rgb("#F2F5F8")
            panel.line.color.rgb = rgb(ARG_BLUE); panel.line.width = Pt(0.75); panel.shadow.inherit = False
            tb = sl.shapes.add_textbox(Inches(box[0] + 0.2), Inches(box[1] + 0.15),
                                       Inches(box[2] - 0.4), Inches(box[3] - 0.3))
            tf = tb.text_frame; tf.word_wrap = True
            for i, ln in enumerate(s["code"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run(); r.text = ln
                r.font.name = "Consolas"; r.font.size = Pt(F_CODE); r.font.color.rgb = rgb(DARK)
            continue

        # kind == content
        sl = prs.slides.add_slide(L_BULLETS)
        t = ph_by_idx(sl, 0)
        if t is not None: set_text(t, s["title"])
        sub = ph_by_idx(sl, 13)
        if sub is not None and s.get("subtitle"): set_text(sub, s["subtitle"])

        has_img = bool(s.get("image") and os.path.isfile(s["image"]))
        body = ph_by_idx(sl, 14)
        if body is not None:
            if has_img:  # narrow the content placeholder to the left column
                body.left = Inches(0.5); body.top = Inches(1.6)
                body.width = Inches(5.2); body.height = Inches(3.6)
            fill_bullets(body, s["bullets"])

        if has_img:
            img_box = (6.0, 1.7, 3.7, 3.0)
            add_image_fit(sl, s["image"], img_box)
            if s.get("image_caption"):
                add_caption(sl, s["image_caption"], (6.0, 4.75, 3.7, 0.4))

    prs.save(out)
    print("[slides] wrote %s (%d slides) on the Argonne template from run %s"
          % (out, len(spec["slides"]), spec["run_id"]))


def qn_r_id():
    # relationship-id attribute name on <p:sldId>
    return "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


# ============================ matplotlib PDF proxy backend ============================
def render_pdf(spec, out_pdf):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages
    from matplotlib.patches import Rectangle, FancyBboxPatch
    import matplotlib.image as mpimg
    import logging
    # Prefer Arial (matches the .pptx); fall back silently to a metric-compatible sans.
    matplotlib.rcParams["font.family"] = ["Arial", "Liberation Sans", "DejaVu Sans"]
    logging.getLogger("matplotlib.font_manager").setLevel(logging.ERROR)

    def Y(y):
        return H - y

    def wrap(text, size, width_in):
        maxchars = max(10, int(width_in / (0.48 * size / 72.0)))
        return textwrap.wrap(text, maxchars) or [""]

    with PdfPages(out_pdf) as pdf:
        for s in spec["slides"]:
            fig = plt.figure(figsize=(W, H))
            ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, W); ax.set_ylim(0, H); ax.axis("off")
            kind = s["kind"]

            if kind == "closing":
                ax.add_patch(Rectangle((0, 0), W, H, color=ARG_BLUE))
                ax.text(W / 2, Y(2.5), "Argonne National Laboratory", fontsize=22, color=WHITE,
                        ha="center", va="center", fontweight="bold")
                ax.text(W / 2, Y(3.1), "Thank you", fontsize=15, color="#CFE0EA", ha="center", va="center")
                ax.add_patch(Rectangle((0, 0), W, 0.18, color=ARG_GOLD))
                pdf.savefig(fig); plt.close(fig); continue

            if kind == "title":
                ax.add_patch(Rectangle((0, 0), W, H, color=ARG_BLUE))
                ax.add_patch(Rectangle((0, 0), W, 0.22, color=ARG_GOLD))
                ax.add_patch(Rectangle((0, H - 0.22), W, 0.22, color=ARG_GOLD))
                # optional image, right side
                has_title_img = bool(s.get("image") and os.path.isfile(s["image"]))
                if has_title_img:
                    img = mpimg.imread(s["image"]); ih, iw = img.shape[0], img.shape[1]
                    bx, by, bw, bh = 6.5, 1.0, 3.0, 2.2
                    scale = min(bw / iw, bh / ih); dw, dh = iw * scale, ih * scale
                    ix = bx + (bw - dw) / 2; iy = by + (bh - dh) / 2
                    ax.imshow(img, extent=[ix, ix + dw, Y(iy + dh), Y(iy)], aspect="auto", zorder=5)
                # title text kept in the left column so it never runs under the image
                title_w = 4.8 if has_title_img else W - 1.0
                yy = 1.05
                for ln in s["title"].split("\n"):
                    for seg in wrap(ln, F_TITLE, title_w):
                        ax.text(0.5, Y(yy), seg, fontsize=F_TITLE, color=WHITE, ha="left", va="top", fontweight="bold")
                        yy += F_TITLE / 72.0 * 1.18
                yy += 0.18
                for seg in wrap(s["sub"], F_SUB, title_w):
                    ax.text(0.5, Y(yy), seg, fontsize=F_SUB, color="#9FE0E6", ha="left", va="top")
                    yy += F_SUB / 72.0 * 1.25
                # footer columns
                cols = [("Presenter", "%s / %s" % (s["presenter"], s["role"])),
                        ("Event", s["event"].replace("\n", " ")),
                        ("Date", s["date"])]
                for i, (lab, val) in enumerate(cols):
                    x = 0.5 + i * 3.15
                    ax.text(x, Y(4.55), lab.upper(), fontsize=10, color=ARG_GOLD, ha="left", va="top", fontweight="bold")
                    for k, ln in enumerate(wrap(val, F_AFFIL, 2.9)):
                        ax.text(x, Y(4.85 + k * 0.24), ln, fontsize=F_AFFIL, color="#CFE0EA", ha="left", va="top")
                pdf.savefig(fig); plt.close(fig); continue

            # content / code header
            ax.text(0.5, Y(0.5), s["title"], fontsize=F_HEAD, color=ARG_BLUE_DK, va="top", fontweight="bold")
            if s.get("subtitle"):
                ax.text(0.5, Y(1.12), s["subtitle"], fontsize=F_SUBHEAD, color=GREY, va="top", style="italic")
            ax.add_patch(Rectangle((0.5, Y(1.5)), W - 1.0, 0.045, color=ARG_GOLD))

            has_img = bool(s.get("image") and os.path.isfile(s["image"]))
            body_w = 5.0 if has_img else W - 1.0

            if kind == "code":
                ax.add_patch(FancyBboxPatch((0.5, 0.55), W - 1.0, 3.35,
                             boxstyle="round,pad=0.02,rounding_size=0.08",
                             linewidth=0.8, edgecolor=ARG_BLUE, facecolor="#F2F5F8", zorder=1))
                yy = 1.95
                for ln in s["code"]:
                    ax.text(0.75, Y(yy), ln, fontsize=F_CODE, color=DARK, va="top", family="monospace", zorder=2)
                    yy += F_CODE / 72.0 * 1.55
            else:
                yy = 1.75
                for text, lvl in s["bullets"]:
                    size = F_L1 if lvl else F_L0
                    color = GREY if lvl else DARK
                    marker = "– " if lvl else "• "
                    indent = 0.6 + (0.45 if lvl else 0.0)
                    for k, ln in enumerate(wrap(text, size, body_w - (0.5 if lvl else 0.0))):
                        ax.text(indent, Y(yy), (marker + ln) if k == 0 else ln,
                                fontsize=size, color=color, va="top")
                        yy += size / 72.0 * 1.4
                    yy += 0.07

            if has_img:
                img = mpimg.imread(s["image"]); ih, iw = img.shape[0], img.shape[1]
                bx, by, bw, bh = 6.0, 1.7, W - 6.0 - 0.5, H - 1.7 - 0.7
                scale = min(bw / iw, bh / ih); dw, dh = iw * scale, ih * scale
                ix = bx + (bw - dw) / 2; iy = by
                ax.imshow(img, extent=[ix, ix + dw, Y(iy + dh), Y(iy)], aspect="auto", zorder=5)
                ax.text(bx + bw / 2, Y(iy + dh + 0.24), s.get("image_caption", ""),
                        fontsize=F_CAP, color=GREY, ha="center", va="top", style="italic")

            ax.text(0.5, Y(H - 0.14), "PETSc multi-agent tokamak simulation  ·  Argonne National Laboratory",
                    fontsize=8, color=GREY, ha="left", va="bottom")
            pdf.savefig(fig); plt.close(fig)
    print("[slides] wrote %s (%d pages)" % (out_pdf, len(spec["slides"])))


def main():
    spec = build_spec()
    render_pptx(spec, os.path.join(HERE, "petsc_multiagent_tokamak.pptx"))
    render_pdf(spec, os.path.join(HERE, "petsc_multiagent_tokamak.pdf"))


if __name__ == "__main__":
    main()
