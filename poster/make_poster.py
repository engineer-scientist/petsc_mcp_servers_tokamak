#!/usr/bin/env python3
"""
make_poster.py -- build the US-RSE 2026 poster from ONE layout spec, with two backends:
  * an EDITABLE PowerPoint  -> poster/USRSE26_poster.pptx   (native text boxes + shapes)
  * a rendered image + PDF  -> poster/USRSE26_poster.png / .pdf  (matplotlib)

Both back-ends use the SAME positions (inches) and font sizes (points). Points are
absolute (1 pt = 1/72 in), so the preview faithfully shows how big the text is on the
48 in x 36 in poster -- use it to check readability without PowerPoint/LibreOffice.

Content/numbers come from artifacts/<latest>/metrics.json + verification.json; figures
from figures/.

Run with a Python that has python-pptx + matplotlib (the tokamak venv has both):
  /home/sarthak.sharma/tokamak/.venv/bin/python poster/make_poster.py
"""
import os
import json
import textwrap

HERE = os.path.dirname(os.path.abspath(__file__))
PROJECT = os.path.dirname(HERE)
ARTIFACTS = os.path.join(PROJECT, "artifacts")
FIGURES = os.path.join(PROJECT, "figures")

W, H = 48.0, 36.0                      # poster size, inches
BLUE = "#002B5C"; TEAL = "#007D8A"; DARK = "#1F1F1F"; GREY = "#4A4A4A"; WHITE = "#FFFFFF"

# ---- poster-appropriate font sizes (points) ----
F_TITLE, F_SUB, F_AUTH = 66, 34, 27
F_HEAD, F_BODY, F_SMALL, F_EQ, F_CODE, F_CAP = 40, 30, 27, 40, 23, 25
BAND_H = 5.8   # title band height (inches)


def latest_run():
    p = os.path.join(ARTIFACTS, "LATEST")
    if os.path.isfile(p):
        return open(p).read().strip()
    runs = sorted(d for d in os.listdir(ARTIFACTS) if d.startswith("run-"))
    return runs[-1] if runs else None


def load(p):
    return json.load(open(p)) if os.path.isfile(p) else {}


IMAGES = os.path.join(PROJECT, "images")


def _img_aspect(path, default=1.4):
    try:
        from PIL import Image
        w, h = Image.open(path).size
        return w / float(h)
    except Exception:
        try:
            import matplotlib.image as mpimg
            im = mpimg.imread(path)
            return im.shape[1] / float(im.shape[0])
        except Exception:
            return default


def _fig_box_h(path, w, titled=True):
    """Height (in) a figure card needs so its image fills the width with no vertical gap."""
    avail_w = w - 1.0
    img_h = avail_w / _img_aspect(path)
    return (1.5 if titled else 0.35) + img_h + 1.35 + 0.2


def _panel_h(body, w):
    """Estimate the height (in) a text panel needs for its content (header + wrapped lines)."""
    avail = w - 0.6
    hh = 1.25 + 0.30                                   # header band + top padding
    for item in body:
        text, size = item[0], item[1]
        role = item[3] if len(item) > 3 else ""
        if role in ("eq", "code"):
            nlines = 1                                 # rendered on one line, not wrapped
        else:
            cpl = max(10, int(avail / (0.52 * size / 72.0)))
            nlines = max(1, -(-len(text) // cpl))      # ceil division
        hh += nlines * (size / 72.0 * 1.32) + 0.12
    return hh + 0.30                                    # bottom padding


def _layout_column(x, colw, items, ytop, ybot, gut=0.55, gut_cap=1.1):
    """Stack items so the column FILLS [ytop, ybot] with minimal empty space: figures get their
    aspect-correct height and text panels get their content height (so no card has an empty
    bottom); the leftover slack goes first into UNIFORM inter-card spacing (capped), and only any
    remainder is shared into the panels."""
    heights = [_fig_box_h(it["path"], colw, "title" in it) if it["type"] == "fig"
               else _panel_h(it["body"], colw) for it in items]
    n = len(items)
    panel_ix = [i for i, it in enumerate(items) if it["type"] == "panel"]
    avail = ybot - ytop
    slack = avail - (sum(heights) + gut * (n - 1))
    if slack >= 0:
        if n > 1:                                       # widen gutters first (up to the cap)
            gut += max(0.0, min(slack / (n - 1), gut_cap - gut))
            slack = avail - (sum(heights) + gut * (n - 1))
        if slack > 0 and panel_ix:                      # remainder → panels
            for i in panel_ix:
                heights[i] += slack / len(panel_ix)
    else:                                               # overflow: tighten gutters, then panels
        gut = max(0.35, gut + slack / (n - 1)) if n > 1 else gut
        over = (sum(heights) + gut * (n - 1)) - avail
        if over > 0 and panel_ix:
            for i in panel_ix:
                heights[i] -= over / len(panel_ix)
    panels, figs, y = [], [], ytop
    for it, hgt in zip(items, heights):
        rec = dict(x=x, y=y, w=colw, h=hgt)
        if it["type"] == "fig":
            rec["path"] = it["path"]; rec["caption"] = it["caption"]
            if "title" in it:
                rec["title"] = it["title"]
            figs.append(rec)
        else:
            rec["title"] = it["title"]; rec["body"] = it["body"]
            panels.append(rec)
        y += hgt + gut
    return panels, figs


def _logo_rgb(path):
    """Return a path to an RGB(A) version of a logo, converting CMYK/other modes (e.g. the
    SUNY .jpg is CMYK, which PowerPoint/matplotlib render with wrong colors)."""
    from PIL import Image
    im = Image.open(path)
    if im.mode in ("RGB", "RGBA", "L", "P"):
        return path
    cache = os.path.join(HERE, ".logo_cache")
    os.makedirs(cache, exist_ok=True)
    out = os.path.join(cache, os.path.splitext(os.path.basename(path))[0] + "_rgb.png")
    im.convert("RGB").save(out)
    return out


def _place_logos(paths, band_y, band_h, target_h=1.5, gap=2.6):
    """Lay out logos in a horizontal row, centered in the poster width and vertically in the
    footer band. Returns [(rgb_path, x, y, w, h), ...]. Scales down if the row is too wide."""
    from PIL import Image
    ars = [Image.open(p).size[0] / Image.open(p).size[1] for p in paths]
    widths = [target_h * ar for ar in ars]
    total = sum(widths) + gap * (len(paths) - 1)
    if total > W - 2.0:                                 # scale the whole row to fit
        sc = (W - 2.0) / total
        target_h *= sc; widths = [w * sc for w in widths]; gap *= sc; total = W - 2.0
    x = (W - total) / 2.0
    y = band_y + (band_h - target_h) / 2.0
    out = []
    for p, w in zip(paths, widths):
        out.append((_logo_rgb(p), x, y, w, target_h))
        x += w + gap
    return out


def build_spec():
    """Return the poster as data: a title band, panels, figures. Positions in inches.

    Content leads with the milestone-9 result: one agent-generated solver verified across
    THREE real-machine shaped equilibria (ITER, NSTX-like, X-point) at 2nd order, with a
    q-profile cross-checked against FreeGS."""
    run_id = latest_run()                                   # the shaped run (artifacts/LATEST)
    run_dir = os.path.join(ARTIFACTS, run_id) if run_id else ""
    m = load(os.path.join(run_dir, "metrics.json"))
    summ = load(os.path.join(run_dir, "shaped_summary.json")).get("machines", {})
    e = m.get("efficiency", {}); h = m.get("human_effort", {})
    lines_gen = h.get("solver_lines_agent_generated") or 252
    wall = e.get("wallclock_seconds_total", 311)
    llm = e.get("approx_llm_completions", 27)

    def q95(machine, default):
        d = summ.get(machine, {})
        return "%.1f" % d["q95"] if d.get("q95") else default

    colw, gut = 14.9, 0.55
    x1 = 0.7; x2 = x1 + colw + gut; x3 = x2 + colw + gut
    ytop, ybot = 6.4, 33.1                              # leave a logo footer below the columns

    # ---- column 1: why + method ------------------------------------------------
    col1 = [
        dict(type="panel", title="1  Motivation", body=[
            ("Building correct, fast HPC simulation code demands scarce, implicit expertise. "
             "Can a system of multiple AI agents automate the path from a plain-language idea "
             "to VERIFIED code for a real fusion problem?", F_BODY, DARK),
            ("We answer yes for the tokamak Grad-Shafranov equilibrium — with verification "
             "as a first-class deliverable.", F_BODY, TEAL, "b"),
        ]),
        dict(type="panel", title="2  The problem", body=[
            ("Axisymmetric ideal-MHD force balance for the poloidal flux psi(R,Z):", F_BODY, DARK),
            ("Δ*psi = -μ₀R² p'(psi) - F F'(psi)", F_EQ, TEAL, "eq"),
            ("Sets the flux surfaces, magnetic axis, plasma shape, and safety factor q.", F_BODY, DARK),
            ("Solov'ev profiles admit an EXACT solution → a rigorous verification anchor.", F_BODY, DARK),
        ]),
        dict(type="fig", title="3  The multi-agent system",
             path=os.path.join(IMAGES, "system_of_multiple_AI_agents_to_automate_simulations.png"),
             caption="Fig. 1  A Problem-Definition layer, an Orchestrator agent + shared/persistent "
                     "memory (Workflow Control), and specialist Modeling / Numerical-Analysis / "
                     "Code-Generation / Visualization agents (Agent Execution)."),
    ]

    # ---- column 2: pipeline + the real-machine equilibria ----------------------
    col2 = [
        dict(type="panel", title="4  The pipeline in action", body=[
            ("Prompt: “the Grad-Shafranov equilibrium for a tokamak plasma.”", F_BODY, DARK, "b"),
            ("Modeling → ‘Grad-Shafranov equation’, steady.", F_BODY, DARK),
            ("Numerical Analysis → nonlinear ⇒ SNES on a DMDA.", F_BODY, DARK),
            ("Code Generation → ONE %d-line parametrized PETSc solver." % lines_gen, F_BODY, DARK),
            ("Compiled & ran on 1 and 4 MPI ranks — no human edits.", F_BODY, TEAL, "b"),
        ]),
        dict(type="panel", title="5  Generated solver (excerpt)", body=[
            ("/* GS operator (normalized x=R/R0), 2nd-order FD */", F_CODE, GREY, "code"),
            ("DMDACreate2d(...,DMDA_STENCIL_STAR,...,&da);", F_CODE, DARK, "code"),
            ("dxx=(u[j][i+1]-2*u[j][i]+u[j][i-1])/hx2;", F_CODE, DARK, "code"),
            ("dyy=(u[j+1][i]-2*u[j][i]+u[j-1][i])/hy2;", F_CODE, DARK, "code"),
            ("f[j][i]=dxx-dx/x+dyy-((1-A)*x*x+A);", F_CODE, DARK, "code"),
            ("/* Dirichlet BC = analytic CF psi (nonzero) */", F_CODE, GREY, "code"),
            ("f[b]=u[b]-PsiExact(&u,x,y);", F_CODE, DARK, "code"),
            ("SNESSetJacobian(snes,J,J,FormJacobian,&u);", F_CODE, DARK, "code"),
            ("SNESSolve(snes,NULL,X);", F_CODE, DARK, "code"),
        ]),
        dict(type="fig", title="6  Real-machine equilibria",
             path=os.path.join(FIGURES, "shaped_equilibria.png"),
             caption="Fig. 2  One agent-generated solver → ITER, spherical-tokamak (NSTX-like) "
                     "and diverted double-null (X-point) equilibria. Red = separatrix; × = X-points."),
    ]

    # ---- column 3: verification, q-profile, results ----------------------------
    col3 = [
        dict(type="fig", title="7  Verification: 2nd order",
             path=os.path.join(FIGURES, "shaped_convergence.png"),
             caption="Fig. 3  Method of manufactured solutions vs the exact Cerfon-Freidberg "
                     "solution: observed order p = 2.00 for all three machines."),
        dict(type="panel", title="8  Safety factor & FreeGS cross-check", body=[
            ("Safety-factor q(ψ_N) computed from the verified flux surfaces:", F_BODY, DARK),
            ("ITER q95≈%s · NSTX q95≈%s · X-point q95≈%s"
             % (q95("iter", "2.9"), q95("nstx", "12"), q95("xpoint", "3.2")), F_BODY, TEAL, "b"),
            ("Cross-checked vs FreeGS’s independent q on the SAME field → agree to < 0.2%.", F_BODY, DARK),
            ("Measured κ, δ match the input shaping to ~1–2%.", F_BODY, DARK),
        ]),
        dict(type="panel", title="9  Results & conclusions", body=[
            ("One %d-line solver, 0 hand-written, verified on 3 real-machine equilibria."
             % lines_gen, F_BODY, DARK),
            ("2nd-order accurate (p = 2.00) on every case; runs on 1 & 4 MPI ranks.", F_BODY, DARK),
            ("Efficiency: ~%s s wall-clock; ~%s LLM completions." % (wall, llm), F_BODY, DARK),
            ("Verification-driven: exact-solution convergence + independent q cross-check.",
             F_BODY, DARK),
            ("github.com/engineer-scientist/petsc_mcp_servers_tokamak", F_SMALL, TEAL),
        ]),
    ]

    panels, figs = [], []
    for x, col in ((x1, col1), (x2, col2), (x3, col3)):
        pp, ff = _layout_column(x, colw, col, ytop, ybot, gut)
        panels += pp; figs += ff

    title = dict(
        title="A System of Multiple AI Agents for Automating Tokamak-Plasma Simulation for Nuclear Fusion Energy",
        sub="A hierarchical multi-agent PETSc system: from a plain-language prompt to verified, real-machine Grad-Shafranov equilibria",
        auth="Sarthak Sharma (State University of New York at Buffalo)  ·  Dr Junchao Zhang (Mathematics and Computer Science Division, Argonne National Laboratory)  ·  US-RSE 2026",
    )
    footer = dict(y=33.5, h=1.9, logos=[
        os.path.join(IMAGES, "ANL_logo.png"),
        os.path.join(IMAGES, "PETSc_TAO_logo.png"),
        os.path.join(IMAGES, "SUNY_Buffalo_logo.jpg"),
    ])
    return dict(title=title, panels=panels, figs=figs, footer=footer, run_id=run_id)


# ============================ matplotlib preview backend ============================
def render_preview(spec, out_png, out_pdf):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch, Rectangle
    import matplotlib.image as mpimg

    fig = plt.figure(figsize=(W, H))
    ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, W); ax.set_ylim(0, H); ax.axis("off")

    def Y(y_top):
        return H - y_top

    def wrap(text, size, width_in):
        # avg glyph ~0.52 em; em = size/72 in
        maxchars = max(8, int(width_in / (0.52 * size / 72.0)))
        return textwrap.wrap(text, maxchars) or [""]

    def draw_lines(x, y_top, w, body, pad=0.3):
        yy = y_top
        for spec_line in body:
            text, size, color = spec_line[0], spec_line[1], spec_line[2]
            role = spec_line[3] if len(spec_line) > 3 else ""
            bold = role in ("b", "eq")
            mono = role == "code"
            align = "center" if role == "eq" else "left"
            xx = x + w / 2 if align == "center" else x + pad
            fam = "monospace" if mono else "sans-serif"
            for ln in ([text] if (mono or role == "eq") else wrap(text, size, w - 2 * pad)):
                ax.text(xx, Y(yy), ln, fontsize=size, color=color, va="top",
                        ha=align, fontweight="bold" if bold else "normal", family=fam)
                yy += size / 72.0 * 1.32
            yy += 0.12
        return yy

    # title band (title wraps to ~2 lines; then subtitle, then authors)
    ax.add_patch(Rectangle((0, Y(BAND_H)), W, BAND_H, color=BLUE, zorder=0))
    yy = 0.45
    for ln in textwrap.wrap(spec["title"]["title"], 47):
        ax.text(0.8, Y(yy), ln, fontsize=F_TITLE, color=WHITE, va="top", fontweight="bold")
        yy += F_TITLE / 72.0 * 1.16
    yy += 0.18
    ax.text(0.8, Y(yy), spec["title"]["sub"], fontsize=F_SUB, color="#CFE6EA", va="top")
    yy += F_SUB / 72.0 * 1.3 + 0.15
    ax.text(0.8, Y(yy), spec["title"]["auth"], fontsize=F_AUTH, color=WHITE, va="top")

    # panels
    for p in spec["panels"]:
        ax.add_patch(FancyBboxPatch((p["x"], Y(p["y"] + p["h"])), p["w"], p["h"],
                     boxstyle="round,pad=0,rounding_size=0.12", ec=TEAL, fc=WHITE, lw=2))
        ax.add_patch(Rectangle((p["x"], Y(p["y"] + 1.25)), p["w"], 1.25, color=BLUE))
        ax.text(p["x"] + 0.3, Y(p["y"] + 0.28), p["title"], fontsize=F_HEAD, color=WHITE,
                va="top", fontweight="bold")
        draw_lines(p["x"], p["y"] + 1.55, p["w"], p["body"])

    # figures
    for f in spec["figs"]:
        ax.add_patch(FancyBboxPatch((f["x"], Y(f["y"] + f["h"])), f["w"], f["h"],
                     boxstyle="round,pad=0,rounding_size=0.12", ec=TEAL, fc=WHITE, lw=2))
        y_img = f["y"] + 0.3
        if f.get("title"):
            ax.add_patch(Rectangle((f["x"], Y(f["y"] + 1.25)), f["w"], 1.25, color=BLUE))
            ax.text(f["x"] + 0.3, Y(f["y"] + 0.28), f["title"], fontsize=F_HEAD, color=WHITE,
                    va="top", fontweight="bold")
            y_img = f["y"] + 1.5
        cap_h = 1.3
        if f.get("path") and os.path.isfile(f["path"]):
            img = mpimg.imread(f["path"])
            ih, iw = img.shape[0], img.shape[1]
            avail_w = f["w"] - 1.0
            avail_h = f["y"] + f["h"] - y_img - cap_h
            scale = min(avail_w / iw, avail_h / ih) * 72  # px->in via 72 assumption
            dw, dh = iw * scale / 72, ih * scale / 72
            ix = f["x"] + (f["w"] - dw) / 2
            iy = y_img
            ax.imshow(img, extent=[ix, ix + dw, Y(iy + dh), Y(iy)], zorder=5, aspect="auto")
        for i, ln in enumerate(wrap(f["caption"], F_CAP, f["w"] - 0.8)):
            ax.text(f["x"] + f["w"] / 2, Y(f["y"] + f["h"] - cap_h + 0.15 + i * F_CAP / 72 * 1.3),
                    ln, fontsize=F_CAP, color=GREY, va="top", ha="center")

    # footer: separator rule + institution logos
    ft = spec.get("footer")
    if ft:
        ax.add_patch(Rectangle((0.7, Y(ft["y"] - 0.25)), W - 1.4, 0.03, color=TEAL, zorder=6))
        for path, x, y, w, hh in _place_logos(ft["logos"], ft["y"], ft["h"]):
            img = mpimg.imread(path)
            ax.imshow(img, extent=[x, x + w, Y(y + hh), Y(y)], zorder=7, aspect="auto")

    fig.savefig(out_png, dpi=64); fig.savefig(out_pdf); plt.close(fig)
    print("[poster] wrote %s and %s" % (out_png, out_pdf))


# ============================ python-pptx backend ============================
def render_pptx(spec, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    def rgb(hexs):
        return RGBColor(int(hexs[1:3], 16), int(hexs[3:5], 16), int(hexs[5:7], 16))

    prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def box(x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, (text, size, color, *rest) in enumerate(lines):
            role = rest[0] if rest else ""
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER if role == "eq" else PP_ALIGN.LEFT
            p.space_after = Pt(6)
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = rgb(color)
            r.font.bold = role in ("b", "eq")
            if role == "code":
                r.font.name = "Consolas"
        return tb

    def rect(x, y, w, h, color, line=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(color)
        if line:
            s.line.color.rgb = rgb(line); s.line.width = Pt(1.5)
        else:
            s.line.fill.background()
        s.shadow.inherit = False
        return s

    def card(x, y, w, h):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(WHITE)
        s.line.color.rgb = rgb(TEAL); s.line.width = Pt(2); s.shadow.inherit = False
        return s

    def header(x, y, w, title):
        hb = rect(x, y, w, 1.25, BLUE)
        tf = hb.text_frame; tf.margin_left = Inches(0.3); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = tf.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(F_HEAD); r.font.bold = True; r.font.color.rgb = rgb(WHITE)

    # title band
    rect(0, 0, W, BAND_H, BLUE)
    box(0.8, 0.35, W - 1.6, BAND_H - 0.6, [
        (spec["title"]["title"], F_TITLE, WHITE, "b"),
        (spec["title"]["sub"], F_SUB, "#CFE6EA"),
        (spec["title"]["auth"], F_AUTH, WHITE),
    ])

    for p in spec["panels"]:
        card(p["x"], p["y"], p["w"], p["h"]); header(p["x"], p["y"], p["w"], p["title"])
        box(p["x"] + 0.3, p["y"] + 1.45, p["w"] - 0.6, p["h"] - 1.7, p["body"])

    for f in spec["figs"]:
        card(f["x"], f["y"], f["w"], f["h"])
        y_img = f["y"] + 0.35
        if f.get("title"):
            header(f["x"], f["y"], f["w"], f["title"]); y_img = f["y"] + 1.5
        if f.get("path") and os.path.isfile(f["path"]):
            from PIL import Image
            iw, ih = Image.open(f["path"]).size
            cap_h = 1.35
            avail_w = f["w"] - 1.0; avail_h = f["y"] + f["h"] - y_img - cap_h
            scale = min(avail_w / iw, avail_h / ih)
            dw, dh = iw * scale, ih * scale
            slide.shapes.add_picture(f["path"], Inches(f["x"] + (f["w"] - dw) / 2),
                                     Inches(y_img), height=Inches(dh))
        box(f["x"] + 0.3, f["y"] + f["h"] - 1.35, f["w"] - 0.6, 1.2,
            [(f["caption"], F_CAP, GREY)], anchor=MSO_ANCHOR.TOP)

    # footer: separator rule + institution logos
    ft = spec.get("footer")
    if ft:
        rect(0.7, ft["y"] - 0.25, W - 1.4, 0.03, TEAL)
        for path, x, y, w, hh in _place_logos(ft["logos"], ft["y"], ft["h"]):
            slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(hh))

    prs.save(out)
    print("[poster] wrote %s (%.0f x %.0f in) from run %s" % (out, W, H, spec["run_id"]))


def main():
    spec = build_spec()
    render_pptx(spec, os.path.join(HERE, "USRSE26_poster.pptx"))
    render_preview(spec, os.path.join(HERE, "USRSE26_poster.png"),
                   os.path.join(HERE, "USRSE26_poster.pdf"))


if __name__ == "__main__":
    main()
